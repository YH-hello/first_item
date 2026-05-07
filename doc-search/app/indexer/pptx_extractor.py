"""
PowerPoint(PPTX) 텍스트 및 슬라이드 정보 추출기
라이브러리: python-pptx
딥링크: win32com PowerPoint.Application → GotoSlide(slide_idx)
"""
import logging
from typing import List

from pptx import Presentation

from .base_extractor import BaseExtractor, PageChunk

logger = logging.getLogger(__name__)


class PPTXExtractor(BaseExtractor):
    """python-pptx를 이용한 슬라이드별 텍스트 및 노트 추출"""

    def extract(self) -> List[PageChunk]:
        chunks: List[PageChunk] = []
        try:
            prs = Presentation(str(self.file_path))
            total = len(prs.slides)

            for slide_idx, slide in enumerate(prs.slides, start=1):
                texts: List[str] = []

                # 슬라이드 내 모든 텍스트 프레임 수집
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            line = " ".join(run.text for run in para.runs).strip()
                            if line:
                                texts.append(line)

                # 발표자 노트
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        texts.append(f"[노트] {notes}")

                text = "\n".join(texts)
                if not text.strip():
                    continue

                slide_id = slide.slide_id
                page_label = f"{slide_idx}번째 슬라이드"

                # win32com을 통한 특정 슬라이드 이동 명령
                abs_path = str(self.file_path.resolve()).replace("\\", "\\\\")
                page_link_cmd = (
                    f"python -c \""
                    f"import win32com.client; "
                    f"ppt=win32com.client.Dispatch('PowerPoint.Application'); "
                    f"ppt.Visible=True; "
                    f"prs=ppt.Presentations.Open(r'{abs_path}'); "
                    f"prs.Windows(1).View.GotoSlide({slide_idx})\""
                )

                chunks.append(
                    self._build_chunk(
                        text=text,
                        page_num=slide_idx,
                        page_label=page_label,
                        page_link_cmd=page_link_cmd,
                        extra_meta={
                            "slide_id": slide_id,
                            "total_slides": total,
                        },
                    )
                )

        except Exception as exc:
            logger.error("[PPTX] 추출 오류 %s: %s", self.file_path, exc, exc_info=True)

        return chunks
